"""
데모 D 폴백 — 같은 데이터로 PPT + HTML + React 3종 출력

메시지: "CLAUDE.md에 규칙 한 번 쓰면 출력 매체는 자유롭게 선택 가능"

1. overseas_sales.xlsx → 집계된 dashboard data
2. PPT 2장 (python-pptx)
3. HTML 대시보드 (Chart.js CDN)
4. React HTML + JSX 컴포넌트 (Babel CDN)

실행:
  py week07-hands-on/demos/dashboard_3format/fallback_build_all.py

출력:
  output/매출보고서.pptx
  output/dashboard.html
  output/dashboard_react.html
  output/DashBoard.jsx
"""

from __future__ import annotations

import json
import sys
from collections import defaultdict
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Inches, Pt
from pptx.dml.color import RGBColor

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8")


HERE = Path(__file__).parent
DATA_PATH = HERE / "data" / "overseas_sales.xlsx"
OUT_DIR = HERE / "output"
OUT_DIR.mkdir(exist_ok=True)

NAVY = RGBColor(0x1F, 0x4E, 0x78)
GRAY = RGBColor(0xD9, 0xD9, 0xD9)
RED = RGBColor(0xC0, 0x00, 0x00)


# ─────────────────────────────────────────────────────────────
# 데이터 로딩 & 집계
# ─────────────────────────────────────────────────────────────


def load_and_aggregate() -> dict:
    wb = load_workbook(DATA_PATH, data_only=True)

    # 법인마스터
    legions: dict[str, dict] = {}
    ws = wb["법인마스터"]
    for row in ws.iter_rows(min_row=2, values_only=True):
        code, name, country, region, currency, contact, email = row
        if not code:
            continue
        legions[code] = {
            "name": name, "country": country, "region": region,
            "currency": currency, "contact": contact, "email": email,
        }

    # 환율
    rates: dict[str, float] = {}
    ws = wb["환율"]
    for row in ws.iter_rows(min_row=2, values_only=True):
        currency, rate, date = row
        if currency:
            rates[currency] = float(rate)

    # 월별매출 집계 (2026년 상반기 전체)
    ws = wb["월별매출"]
    legion_totals: dict[str, dict] = defaultdict(lambda: {"sales_local": 0.0, "cost": 0.0, "sga": 0.0})
    for row in ws.iter_rows(min_row=2, values_only=True):
        code, name, month, pg, sales, cur, cost_rate, cost, sga_rate, sga, op = row
        if not code:
            continue
        legion_totals[code]["sales_local"] += float(sales or 0)
        legion_totals[code]["cost"] += float(cost or 0)
        legion_totals[code]["sga"] += float(sga or 0)

    # KRW 환산 + 영업이익률
    legion_summary = []
    total_revenue_krw = 0.0
    total_cost_krw = 0.0
    total_op_krw = 0.0
    region_totals = defaultdict(lambda: {"revenue_krw": 0.0, "legion_count": 0})

    for code, totals in legion_totals.items():
        info = legions[code]
        rate = rates.get(info["currency"], 1.0)
        revenue_krw = totals["sales_local"] * rate
        cost_krw = totals["cost"] * rate
        sga_krw = totals["sga"] * rate
        op_krw = revenue_krw - cost_krw - sga_krw
        op_margin = (op_krw / revenue_krw * 100) if revenue_krw else 0

        legion_summary.append({
            "code": code,
            "name": info["name"],
            "country": info["country"],
            "region": info["region"],
            "currency": info["currency"],
            "contact": info["contact"],
            "sales_local": totals["sales_local"],
            "revenue_krw": revenue_krw,
            "cost_krw": cost_krw,
            "sga_krw": sga_krw,
            "op_krw": op_krw,
            "op_margin": op_margin,
        })

        total_revenue_krw += revenue_krw
        total_cost_krw += cost_krw
        total_op_krw += op_krw
        region_totals[info["region"]]["revenue_krw"] += revenue_krw
        region_totals[info["region"]]["legion_count"] += 1

    # 상위 10개 법인 (매출 기준)
    top10 = sorted(legion_summary, key=lambda x: -x["revenue_krw"])[:10]

    return {
        "total_revenue_krw": total_revenue_krw,
        "total_cost_krw": total_cost_krw,
        "total_op_krw": total_op_krw,
        "overall_op_margin": (total_op_krw / total_revenue_krw * 100) if total_revenue_krw else 0,
        "region_totals": dict(region_totals),
        "top10": top10,
        "all_legions": legion_summary,
    }


def to_trillion_won(amount: float) -> str:
    """원 → 조/억 표기"""
    if abs(amount) >= 1_000_000_000_000:
        return f"{amount / 1_000_000_000_000:.2f}조원"
    if abs(amount) >= 100_000_000:
        return f"{amount / 100_000_000:.0f}억원"
    return f"{amount:,.0f}원"


# ─────────────────────────────────────────────────────────────
# 출력 1: PowerPoint
# ─────────────────────────────────────────────────────────────


def build_ppt(data: dict) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 와이드
    prs.slide_height = Inches(7.5)

    # ── Slide 1: 전체 요약 ──
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "해외법인 매출 현황 (2026년 상반기)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # KPI 박스 3개 (총매출/총원가/총영업이익)
    kpi_data = [
        ("총매출", to_trillion_won(data["total_revenue_krw"]), NAVY),
        ("총원가", to_trillion_won(data["total_cost_krw"]), RGBColor(0x70, 0x70, 0x70)),
        ("총영업이익", to_trillion_won(data["total_op_krw"]) +
                    f"\n(이익률 {data['overall_op_margin']:.1f}%)", RED),
    ]
    kpi_y = Inches(1.3)
    kpi_w = Inches(4.0)
    kpi_h = Inches(1.5)
    gap = Inches(0.15)
    for i, (label, value, color) in enumerate(kpi_data):
        x = Inches(0.5) + i * (kpi_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, kpi_y, kpi_w, kpi_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(14)
        p1.font.color.rgb = RGBColor(0x60, 0x60, 0x60)
        p1.alignment = 2  # center
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = color
        p2.alignment = 2

    # 지역별 매출 바 차트
    chart_data = CategoryChartData()
    regions = ["Americas", "APAC", "EMEA", "Oceania"]
    chart_data.categories = regions
    chart_data.add_series(
        "매출 (조원)",
        [data["region_totals"].get(r, {}).get("revenue_krw", 0) / 1_000_000_000_000 for r in regions],
    )
    cx, cy, cw, ch = Inches(0.5), Inches(3.2), Inches(12.3), Inches(4.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, cx, cy, cw, ch, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "지역별 매출 규모 (조원)"
    chart.has_legend = False

    # ── Slide 2: 상위 10개 법인 ──
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "매출 상위 10개 법인"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 테이블
    headers = ["순위", "법인명", "국가", "매출(원화)", "영업이익률", "담당자"]
    rows = len(data["top10"]) + 1
    cols = len(headers)
    tx, ty, tw, th = Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5)
    table_shape = slide.shapes.add_table(rows, cols, tx, ty, tw, th)
    table = table_shape.table

    col_widths = [Inches(0.8), Inches(3.2), Inches(1.8), Inches(2.2), Inches(2.0), Inches(2.3)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # 헤더
    for c_idx, h in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            para.font.bold = True
            para.font.size = Pt(14)
            para.alignment = 2

    # 데이터 행
    for idx, legion in enumerate(data["top10"], start=1):
        row = [
            str(idx),
            legion["name"],
            legion["country"],
            to_trillion_won(legion["revenue_krw"]),
            f"{legion['op_margin']:.1f}%",
            legion["contact"],
        ]
        is_highlight = legion["op_margin"] >= 20  # 20% 이상 강조
        for c_idx, v in enumerate(row):
            cell = table.cell(idx, c_idx)
            cell.text = v
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                if is_highlight and c_idx == 4:  # 영업이익률 컬럼만 빨강
                    para.font.color.rgb = RED
                    para.font.bold = True

    out_path = OUT_DIR / "매출보고서.pptx"
    prs.save(out_path)
    return out_path


# ─────────────────────────────────────────────────────────────
# 출력 2: HTML 대시보드 (Chart.js)
# ─────────────────────────────────────────────────────────────


def build_html(data: dict) -> Path:
    out_path = OUT_DIR / "dashboard.html"

    # 데이터를 JS 객체로 직렬화
    region_labels = list(data["region_totals"].keys())
    region_values = [data["region_totals"][r]["revenue_krw"] / 1_000_000_000_000 for r in region_labels]
    legions_js = [
        {
            "code": l["code"],
            "name": l["name"],
            "country": l["country"],
            "region": l["region"],
            "revenue_krw": l["revenue_krw"],
            "op_margin": l["op_margin"],
            "contact": l["contact"],
        }
        for l in data["all_legions"]
    ]

    html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<title>해외법인 매출 현황 (2026년 상반기)</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.1/dist/chart.umd.min.js"></script>
<style>
  body {{ font-family: 'Malgun Gothic', sans-serif; margin: 0; padding: 24px; background: #f5f5f5; color: #222; }}
  h1 {{ color: #1F4E78; margin-bottom: 24px; }}
  .kpi-row {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 16px; margin-bottom: 24px; }}
  .kpi {{ background: white; border-left: 4px solid #1F4E78; padding: 20px; border-radius: 8px;
          box-shadow: 0 2px 6px rgba(0,0,0,.05); }}
  .kpi .label {{ color: #666; font-size: 14px; }}
  .kpi .value {{ font-size: 28px; font-weight: bold; color: #1F4E78; margin-top: 4px; }}
  .kpi.cost {{ border-left-color: #707070; }}
  .kpi.cost .value {{ color: #707070; }}
  .kpi.profit {{ border-left-color: #C00000; }}
  .kpi.profit .value {{ color: #C00000; }}
  .controls {{ background: white; padding: 16px; border-radius: 8px; margin-bottom: 16px; }}
  .controls label {{ margin-right: 12px; font-weight: bold; }}
  select {{ padding: 6px 12px; font-size: 14px; border-radius: 4px; border: 1px solid #ccc; }}
  .chart-box {{ background: white; padding: 20px; border-radius: 8px; margin-bottom: 16px; height: 360px; }}
  table {{ width: 100%; background: white; border-collapse: collapse; border-radius: 8px; overflow: hidden; }}
  th {{ background: #1F4E78; color: white; padding: 12px; text-align: left; }}
  td {{ padding: 10px 12px; border-bottom: 1px solid #eee; }}
  tr:hover {{ background: #f9f9f9; }}
  .margin-high {{ color: #C00000; font-weight: bold; }}
  .right {{ text-align: right; }}
</style>
</head>
<body>
<h1>해외법인 매출 현황 (2026년 상반기)</h1>

<div class="kpi-row">
  <div class="kpi"><div class="label">총매출</div><div class="value">{to_trillion_won(data["total_revenue_krw"])}</div></div>
  <div class="kpi cost"><div class="label">총원가</div><div class="value">{to_trillion_won(data["total_cost_krw"])}</div></div>
  <div class="kpi profit"><div class="label">총영업이익 (이익률 {data["overall_op_margin"]:.1f}%)</div>
       <div class="value">{to_trillion_won(data["total_op_krw"])}</div></div>
</div>

<div class="controls">
  <label for="regionFilter">지역 필터:</label>
  <select id="regionFilter">
    <option value="">전체</option>
    <option>Americas</option>
    <option>APAC</option>
    <option>EMEA</option>
    <option>Oceania</option>
  </select>
</div>

<div class="chart-box"><canvas id="regionChart"></canvas></div>

<h2>법인별 매출 상세</h2>
<table id="legionTable">
  <thead><tr>
    <th>법인코드</th><th>법인명</th><th>국가</th><th>지역</th>
    <th class="right">매출 (원화)</th><th class="right">영업이익률</th><th>담당자</th>
  </tr></thead>
  <tbody></tbody>
</table>

<script>
const REGION_DATA = {{ labels: {json.dumps(region_labels)}, values: {json.dumps(region_values)} }};
const LEGIONS = {json.dumps(legions_js, ensure_ascii=False)};

function formatWon(v) {{
  if (Math.abs(v) >= 1e12) return (v / 1e12).toFixed(2) + '조원';
  if (Math.abs(v) >= 1e8)  return (v / 1e8).toFixed(0) + '억원';
  return Math.round(v).toLocaleString() + '원';
}}

const ctx = document.getElementById('regionChart').getContext('2d');
const regionChart = new Chart(ctx, {{
  type: 'bar',
  data: {{
    labels: REGION_DATA.labels,
    datasets: [{{ label: '매출 (조원)', data: REGION_DATA.values, backgroundColor: '#1F4E78' }}]
  }},
  options: {{
    responsive: true, maintainAspectRatio: false,
    plugins: {{ legend: {{ display: false }}, title: {{ display: true, text: '지역별 매출 규모 (조원)' }} }}
  }}
}});

function renderTable(region) {{
  const tbody = document.querySelector('#legionTable tbody');
  tbody.innerHTML = '';
  const filtered = region ? LEGIONS.filter(l => l.region === region) : LEGIONS;
  filtered.sort((a, b) => b.revenue_krw - a.revenue_krw);
  for (const l of filtered) {{
    const marginClass = l.op_margin >= 20 ? 'margin-high' : '';
    tbody.innerHTML += `
      <tr>
        <td>${{l.code}}</td>
        <td>${{l.name}}</td>
        <td>${{l.country}}</td>
        <td>${{l.region}}</td>
        <td class="right">${{formatWon(l.revenue_krw)}}</td>
        <td class="right ${{marginClass}}">${{l.op_margin.toFixed(1)}}%</td>
        <td>${{l.contact}}</td>
      </tr>`;
  }}
}}

document.getElementById('regionFilter').addEventListener('change', e => renderTable(e.target.value));
renderTable('');
</script>
</body>
</html>
"""
    out_path.write_text(html, encoding="utf-8")
    return out_path


# ─────────────────────────────────────────────────────────────
# 출력 3: React (Babel inline) + 순수 JSX 파일
# ─────────────────────────────────────────────────────────────


def build_react(data: dict) -> tuple[Path, Path]:
    out_html = OUT_DIR / "dashboard_react.html"
    out_jsx = OUT_DIR / "DashBoard.jsx"

    region_labels = list(data["region_totals"].keys())
    region_values = [data["region_totals"][r]["revenue_krw"] / 1_000_000_000_000 for r in region_labels]
    legions_js = [
        {
            "code": l["code"], "name": l["name"], "country": l["country"], "region": l["region"],
            "revenue_krw": l["revenue_krw"], "op_margin": l["op_margin"], "contact": l["contact"],
        }
        for l in data["all_legions"]
    ]

    # 순수 JSX 컴포넌트 (사내 프로젝트 이식용)
    jsx = """// DashBoard.jsx
// 사내 React 웹앱에 그대로 복붙해서 쓸 수 있는 독립 컴포넌트.
// Chart.js react wrapper(react-chartjs-2)를 쓰면 차트도 native로 가능.

import React, { useState, useMemo } from 'react';

const LEGIONS = __LEGIONS__;
const REGION_DATA = __REGION_DATA__;
const TOTALS = __TOTALS__;

function formatWon(v) {
  if (Math.abs(v) >= 1e12) return (v / 1e12).toFixed(2) + '조원';
  if (Math.abs(v) >= 1e8)  return (v / 1e8).toFixed(0) + '억원';
  return Math.round(v).toLocaleString() + '원';
}

export default function DashBoard() {
  const [region, setRegion] = useState('');

  const filtered = useMemo(() => {
    const base = region ? LEGIONS.filter(l => l.region === region) : LEGIONS;
    return [...base].sort((a, b) => b.revenue_krw - a.revenue_krw);
  }, [region]);

  return (
    <div style={{ fontFamily: 'Malgun Gothic, sans-serif', padding: 24, background: '#f5f5f5' }}>
      <h1 style={{ color: '#1F4E78' }}>해외법인 매출 현황 (2026년 상반기)</h1>
      <div style={{ display: 'grid', gridTemplateColumns: 'repeat(3, 1fr)', gap: 16, marginBottom: 24 }}>
        <KPI label="총매출" value={formatWon(TOTALS.revenue)} color="#1F4E78" />
        <KPI label="총원가" value={formatWon(TOTALS.cost)} color="#707070" />
        <KPI label={`총영업이익 (이익률 ${TOTALS.margin.toFixed(1)}%)`}
             value={formatWon(TOTALS.op)} color="#C00000" />
      </div>
      <div style={{ background: 'white', padding: 16, borderRadius: 8, marginBottom: 16 }}>
        <label style={{ marginRight: 12, fontWeight: 'bold' }}>지역 필터:</label>
        <select value={region} onChange={e => setRegion(e.target.value)}>
          <option value="">전체</option>
          {['Americas', 'APAC', 'EMEA', 'Oceania'].map(r =>
            <option key={r} value={r}>{r}</option>
          )}
        </select>
      </div>
      <table style={{ width: '100%', background: 'white', borderCollapse: 'collapse', borderRadius: 8, overflow: 'hidden' }}>
        <thead><tr style={{ background: '#1F4E78', color: 'white' }}>
          {['법인코드', '법인명', '국가', '지역', '매출 (원화)', '영업이익률', '담당자'].map(h =>
            <th key={h} style={{ padding: 12, textAlign: 'left' }}>{h}</th>
          )}
        </tr></thead>
        <tbody>
          {filtered.map(l =>
            <tr key={l.code} style={{ borderBottom: '1px solid #eee' }}>
              <td style={{ padding: 10 }}>{l.code}</td>
              <td style={{ padding: 10 }}>{l.name}</td>
              <td style={{ padding: 10 }}>{l.country}</td>
              <td style={{ padding: 10 }}>{l.region}</td>
              <td style={{ padding: 10, textAlign: 'right' }}>{formatWon(l.revenue_krw)}</td>
              <td style={{ padding: 10, textAlign: 'right',
                           color: l.op_margin >= 20 ? '#C00000' : 'inherit',
                           fontWeight: l.op_margin >= 20 ? 'bold' : 'normal' }}>
                {l.op_margin.toFixed(1)}%
              </td>
              <td style={{ padding: 10 }}>{l.contact}</td>
            </tr>
          )}
        </tbody>
      </table>
    </div>
  );
}

function KPI({ label, value, color }) {
  return (
    <div style={{ background: 'white', borderLeft: `4px solid ${color}`, padding: 20, borderRadius: 8 }}>
      <div style={{ color: '#666', fontSize: 14 }}>{label}</div>
      <div style={{ fontSize: 28, fontWeight: 'bold', color, marginTop: 4 }}>{value}</div>
    </div>
  );
}
"""

    totals = {
        "revenue": data["total_revenue_krw"],
        "cost": data["total_cost_krw"],
        "op": data["total_op_krw"],
        "margin": data["overall_op_margin"],
    }
    jsx = (jsx
           .replace("__LEGIONS__", json.dumps(legions_js, ensure_ascii=False))
           .replace("__REGION_DATA__", json.dumps({"labels": region_labels, "values": region_values}))
           .replace("__TOTALS__", json.dumps(totals)))

    out_jsx.write_text(jsx, encoding="utf-8")

    # 단독 실행 HTML (React + Babel inline)
    html = """<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<title>해외법인 매출 현황 (React)</title>
<script src="https://unpkg.com/react@18/umd/react.production.min.js" crossorigin></script>
<script src="https://unpkg.com/react-dom@18/umd/react-dom.production.min.js" crossorigin></script>
<script src="https://unpkg.com/@babel/standalone/babel.min.js"></script>
<style>body { margin: 0; }</style>
</head>
<body>
<div id="root"></div>
<script type="text/babel" data-type="module">
"""
    # JSX 컴포넌트 본체를 Babel 태그 안에 삽입 (import/export 제거)
    jsx_inline = (jsx
                  .replace("import React, { useState, useMemo } from 'react';", "const { useState, useMemo } = React;")
                  .replace("export default function DashBoard()", "function DashBoard()"))
    html += jsx_inline
    html += """
ReactDOM.createRoot(document.getElementById('root')).render(<DashBoard />);
</script>
</body>
</html>
"""
    out_html.write_text(html, encoding="utf-8")

    return out_html, out_jsx


# ─────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────


def main() -> None:
    if not DATA_PATH.exists():
        print(f"[에러] 데이터 파일 없음: {DATA_PATH}")
        print(f"       prep/make_overseas_sales.py 먼저 실행해주세요.")
        sys.exit(1)

    print(f"[시작] 데이터 로딩 → {DATA_PATH.name}")
    data = load_and_aggregate()
    print(f"  총매출:    {to_trillion_won(data['total_revenue_krw'])}")
    print(f"  총원가:    {to_trillion_won(data['total_cost_krw'])}")
    print(f"  총영업이익: {to_trillion_won(data['total_op_krw'])} (이익률 {data['overall_op_margin']:.1f}%)")
    print(f"  지역 수:   {len(data['region_totals'])}개, 법인 수: {len(data['all_legions'])}개")
    print()

    print(f"[1/3] PPT 생성...")
    ppt_path = build_ppt(data)
    print(f"  [OK] {ppt_path.name}")

    print(f"[2/3] HTML 대시보드 생성...")
    html_path = build_html(data)
    print(f"  [OK] {html_path.name}")

    print(f"[3/3] React HTML + JSX 파일 생성...")
    react_html, jsx_path = build_react(data)
    print(f"  [OK] {react_html.name}")
    print(f"  [OK] {jsx_path.name}")

    print()
    print(f"[완료] 3종 출력 모두 {OUT_DIR}에 저장됨")


if __name__ == "__main__":
    main()
